import re
import os
import json
import pickle
import numpy as np
from datetime import datetime
from typing import List, Dict, Tuple
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.multiclass import OneVsRestClassifier
from sklearn.linear_model import LogisticRegression
from sklearn.model_selection import train_test_split, cross_val_score, StratifiedKFold
from sklearn.metrics import classification_report, confusion_matrix, accuracy_score
import pdfplumber
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
import io
import base64
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd

# Télécharger TOUS les modèles NLTK nécessaires (une fois pour toutes)
nltk.download('stopwords', quiet=True)
nltk.download('punkt', quiet=True)
nltk.download('punkt_tab', quiet=True)
nltk.download('wordnet', quiet=True)

class CVClassifier:
    def __init__(self):
        self.vectorizer = None
        self.model = None
        self.lemmatizer = WordNetLemmatizer()
        self.stop_words = set(stopwords.words('french') + stopwords.words('english'))
        
        # Métriques stockées
        self.accuracy = None
        self.report = None
        self.cm = None
        self.cv_scores = None
        self.nb_samples = None

    def extract_text_from_file(self, file_path: str) -> str:
        filename = os.path.basename(file_path).lower()
        text = ''

        try:
            # PDF: first try pdfplumber (digital PDF), then OCR if empty
            if filename.endswith('.pdf'):
                try:
                    with pdfplumber.open(file_path) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + ' '
                except Exception as e:
                    print(f"pdfplumber failed for {filename}: {e}")

                # si aucun texte obtenu, tenter OCR via pdf2image + pytesseract
                if not text.strip():
                    try:
                        from pdf2image import convert_from_path
                        from PIL import Image
                        import pytesseract

                        images = convert_from_path(file_path, dpi=300)
                        for img in images:
                            ocr_text = pytesseract.image_to_string(img, lang='fra+eng')
                            if ocr_text:
                                text += ocr_text + ' '
                    except Exception as e:
                        print(f"OCR (pdf2image/pytesseract) failed for {filename}: {e}")

            # Images: use pytesseract
            elif filename.endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff')):
                try:
                    from PIL import Image
                    import pytesseract

                    img = Image.open(file_path)
                    ocr_text = pytesseract.image_to_string(img, lang='fra+eng')
                    if ocr_text:
                        text += ocr_text + ' '
                except Exception as e:
                    print(f"OCR image failed for {filename}: {e}")

            elif filename.endswith('.txt'):
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()

            elif filename.endswith('.docx'):
                from docx import Document
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text += para.text + ' '
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            text += cell.text + ' '

            elif filename.endswith('.pptx'):
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + ' '

            elif filename.endswith('.xlsx'):
                import openpyxl
                wb = openpyxl.load_workbook(file_path)
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    for row in sheet.iter_rows(values_only=True):
                        text += ' '.join([str(cell) if cell is not None else '' for cell in row]) + ' '

        except Exception as e:
            print(f"Erreur extraction {filename}: {e}")

        return text.strip()

    def preprocess_text(self, text: str) -> str:
        text = text.lower()
        text = re.sub(r'[^a-zàâäéèêëïîôùûüÿæœç\s]', ' ', text)
        text = re.sub(r'\s+', ' ', text).strip()
        tokens = word_tokenize(text)
        # Simplifier le préprocessing pour les gros datasets - supprimer la lemmatization
        tokens = [word for word in tokens if word not in self.stop_words and len(word) > 2]
        return ' '.join(tokens)

    def load_dataset(self, csv_path: str):
        """Charge un dataset depuis un fichier CSV ou TXT avec colonnes text et label"""
        if not os.path.exists(csv_path):
            raise FileNotFoundError(f"Fichier non trouvé: {csv_path}")
        
        # Essayer de lire comme CSV d'abord
        try:
            if csv_path.endswith('.txt'):
                # Pour les fichiers .txt, essayer de les lire comme CSV
                df = pd.read_csv(csv_path, sep=',', quotechar='"', encoding='utf-8')
            else:
                df = pd.read_csv(csv_path)
            
            # Mapping des colonnes pour différents formats de dataset
            if 'text' in df.columns and 'label' in df.columns:
                # Format standard (jobs_dataset_large.csv)
                text_col, label_col = 'text', 'label'
            elif 'Resume_str' in df.columns and 'Category' in df.columns:
                # Format Resume.csv
                text_col, label_col = 'Resume_str', 'Category'
            elif 'Text' in df.columns and 'Category' in df.columns:
                # Format Dataset.txt
                text_col, label_col = 'Text', 'Category'
            else:
                available_cols = df.columns.tolist()
                raise ValueError(f"Colonnes non reconnues. Disponibles: {available_cols}")
            
            texts = []
            labels = []
            
            for _, row in df.iterrows():
                text = str(row[text_col]).strip()
                label = str(row[label_col]).strip()
                
                if text and label and len(text) > 10:  # Texte suffisamment long
                    texts.append(text)
                    labels.append(label)
                    
        except Exception as e:
            # Si la lecture CSV échoue, essayer une autre approche
            print(f"Lecture CSV échouée: {e}. Tentative de lecture ligne par ligne...")
            
            texts = []
            labels = []
            current_category = None
            current_text = []
            
            with open(csv_path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # Si la ligne contient une virgule, c'est probablement "Category,Text"
                    if ',' in line and not current_category:
                        parts = line.split(',', 1)
                        if len(parts) == 2:
                            current_category = parts[0].strip()
                            text_part = parts[1].strip().strip('"')
                            if text_part:
                                current_text.append(text_part)
                    elif current_category:
                        # Continuer à accumuler le texte pour cette catégorie
                        current_text.append(line)
                    elif line and not current_category:
                        # Nouvelle catégorie
                        current_category = line
                        current_text = []
            
            # Traiter les données accumulées
            if current_category and current_text:
                texts.append(' '.join(current_text))
                labels.append(current_category)
        
        # Vérifier que chaque classe a au moins 1 exemple (modifié pour les petits datasets de test)
        from collections import Counter
        label_counts = Counter(labels)
        min_samples_per_class = min(label_counts.values())
        
        if min_samples_per_class < 1:
            print(f"Attention: Certaines classes n'ont aucun exemple")
            # Garder seulement les classes avec au moins 1 exemple
            valid_labels = [label for label, count in label_counts.items() if count >= 1]
            filtered_texts = []
            filtered_labels = []
            for text, label in zip(texts, labels):
                if label in valid_labels:
                    filtered_texts.append(text)
                    filtered_labels.append(label)
            texts, labels = filtered_texts, filtered_labels
            print(f"Après filtrage: {len(texts)} échantillons et {len(set(labels))} classes")
        
        return texts, labels

    def load_dataset_from_archive(self, archive_path: str, max_samples: int = 1000):
        """Charge un dataset depuis le dossier archive avec les fichiers CSV liés"""
        import os
        
        if not os.path.exists(archive_path):
            raise FileNotFoundError(f"Dossier archive non trouvé: {archive_path}")

        print("Chargement des fichiers CSV...")
        
        # Charger les données de base des personnes (limiter à max_samples)
        people_df = pd.read_csv(os.path.join(archive_path, '01_people.csv')).head(max_samples)
        
        # Créer un ensemble des person_ids valides
        valid_person_ids = set(people_df['person_id'].values)
        
        # Charger et filtrer les compétences seulement pour les personnes valides
        abilities_df = pd.read_csv(os.path.join(archive_path, '02_abilities.csv'))
        abilities_df = abilities_df[abilities_df['person_id'].isin(valid_person_ids)]
        
        # Charger et filtrer l'éducation
        education_df = pd.read_csv(os.path.join(archive_path, '03_education.csv'))
        education_df = education_df[education_df['person_id'].isin(valid_person_ids)]
        
        # Charger et filtrer les expériences
        experience_df = pd.read_csv(os.path.join(archive_path, '04_experience.csv'))
        experience_df = experience_df[experience_df['person_id'].isin(valid_person_ids)]
        
        print("Agrégation des données par personne...")
        
        texts = []
        labels = []
        
        # Grouper les données par person_id pour optimisation
        abilities_grouped = abilities_df.groupby('person_id')['ability'].agg(list).to_dict()
        
        # Pour l'éducation et l'expérience, utiliser une approche différente
        education_grouped = {}
        for person_id, group in education_df.groupby('person_id'):
            education_grouped[person_id] = group.to_dict('records')
            
        experience_grouped = {}
        for person_id, group in experience_df.groupby('person_id'):
            experience_grouped[person_id] = group.to_dict('records')
        
        # Pour chaque personne, créer un texte de CV complet
        for _, person in people_df.iterrows():
            person_id = person['person_id']
            label = str(person['name']).strip() if pd.notna(person['name']) else 'Unknown'
            
            if label == 'Unknown' or label == '' or len(label) < 3:
                continue  # Skip persons without valid labels
                
            cv_text_parts = []
            
            # Ajouter les compétences
            if person_id in abilities_grouped:
                abilities = [str(ability).strip() for ability in abilities_grouped[person_id] if ability and str(ability).strip()]
                if abilities:
                    cv_text_parts.append("COMPETENCES: " + " | ".join(abilities[:20]))  # Limiter à 20 compétences
            
            # Ajouter l'éducation
            if person_id in education_grouped:
                education_text = []
                for edu in education_grouped[person_id][:3]:  # Limiter à 3 formations
                    edu_parts = []
                    if edu.get('institution') and pd.notna(edu['institution']):
                        edu_parts.append(f"Institution: {edu['institution']}")
                    if edu.get('program') and pd.notna(edu['program']):
                        edu_parts.append(f"Program: {edu['program']}")
                    if edu.get('location') and pd.notna(edu['location']):
                        edu_parts.append(f"Location: {edu['location']}")
                    if edu_parts:
                        education_text.append(" | ".join(edu_parts))
                
                if education_text:
                    cv_text_parts.append("EDUCATION: " + " | ".join(education_text))
            
            # Ajouter les expériences
            if person_id in experience_grouped:
                experience_text = []
                for exp in experience_grouped[person_id][:5]:  # Limiter à 5 expériences
                    exp_parts = []
                    if exp.get('title') and pd.notna(exp['title']):
                        exp_parts.append(f"Title: {exp['title']}")
                    if exp.get('firm') and pd.notna(exp['firm']):
                        exp_parts.append(f"Company: {exp['firm']}")
                    if exp.get('location') and pd.notna(exp['location']):
                        exp_parts.append(f"Location: {exp['location']}")
                    if exp_parts:
                        experience_text.append(" | ".join(exp_parts))
                
                if experience_text:
                    cv_text_parts.append("EXPERIENCE: " + " | ".join(experience_text))
            
            # Créer le texte final si on a au moins des compétences ou expériences
            if cv_text_parts:
                full_text = " ".join(cv_text_parts)
                if len(full_text.strip()) > 300:  # Texte suffisamment long
                    texts.append(full_text)
                    labels.append(label)
        
        print(f"Dataset créé avec {len(texts)} échantillons et {len(set(labels))} classes")
        
        # Vérifier que chaque classe a au moins 2 exemples
        from collections import Counter
        label_counts = Counter(labels)
        min_samples_per_class = min(label_counts.values())
        
        if min_samples_per_class < 2:
            print(f"Attention: Certaines classes n'ont qu'un seul exemple (min: {min_samples_per_class})")
            # Filtrer les classes avec moins de 2 exemples
            valid_labels = [label for label, count in label_counts.items() if count >= 2]
            filtered_texts = []
            filtered_labels = []
            for text, label in zip(texts, labels):
                if label in valid_labels:
                    filtered_texts.append(text)
                    filtered_labels.append(label)
            texts, labels = filtered_texts, filtered_labels
            print(f"Après filtrage: {len(texts)} échantillons et {len(set(labels))} classes")
        
        return texts, labels

    def train(self, texts, labels, test_size=0.2):
        if not texts or not labels:
            raise ValueError("Le dataset est vide. Impossible d'entraîner le modèle.")
            
        processed_texts = [self.preprocess_text(text) for text in texts]
        
        # Vérifier si on peut utiliser la stratification
        from collections import Counter
        label_counts = Counter(labels)
        min_samples_per_class = min(label_counts.values())
        nb_classes = len(label_counts)
        
        # Ajuster test_size pour les petits datasets
        min_test_samples = max(nb_classes, 5)  # Au minimum nb_classes ou 5 échantillons
        if len(labels) <= min_test_samples * 2:
            test_size = 0.3  # Utiliser plus d'échantillons pour l'entraînement
            print(f"Dataset petit ({len(labels)} échantillons), ajustement test_size à {test_size}")
        elif len(labels) * test_size < min_test_samples:
            test_size = min_test_samples / len(labels)
            print(f"Ajustement test_size à {test_size:.2f} pour avoir au moins {min_test_samples} échantillons de test")
        
        use_stratification = min_samples_per_class >= 2 and len(labels) * test_size >= nb_classes
        
        if not use_stratification:
            print("Désactivation de la stratification car certaines classes ont moins de 2 exemples ou dataset trop petit")
        
        X_train, X_test, y_train, y_test = train_test_split(
            processed_texts, labels, test_size=test_size, 
            random_state=42, 
            stratify=labels if use_stratification else None
        )
        
        # Ajuster les paramètres TF-IDF selon la taille du dataset
        nb_samples = len(processed_texts)
        if nb_samples < 100:
            # Paramètres pour petits datasets
            tfidf_params = {
                'ngram_range': (1, 2),
                'max_features': min(500, nb_samples * 3),
                'min_df': 1,  # Inclure tous les termes
                'max_df': 0.9
            }
            print(f"Petit dataset ({nb_samples} échantillons), ajustement des paramètres TF-IDF: {tfidf_params}")
        else:
            # Paramètres pour grands datasets
            tfidf_params = {
                'ngram_range': (1, 2),
                'max_features': 300,
                'min_df': 10,
                'max_df': 0.8
            }
        
        self.vectorizer = TfidfVectorizer(**tfidf_params)
        X_train_tfidf = self.vectorizer.fit_transform(X_train)
        X_test_tfidf = self.vectorizer.transform(X_test)
        self.model = OneVsRestClassifier(LogisticRegression(random_state=42, C=10.0, penalty='l2', max_iter=1000))
        self.model.fit(X_train_tfidf, y_train)
        y_pred = self.model.predict(X_test_tfidf)
        self.accuracy = accuracy_score(y_test, y_pred)
        self.report = classification_report(y_test, y_pred, output_dict=True, zero_division=0)
        self.cm = confusion_matrix(y_test, y_pred)
        self.nb_samples = len(texts)

        # Cross-validation pour une meilleure estimation (adaptée aux petits datasets)
        try:
            min_samples_per_class = min(label_counts.values())
            nb_classes = len(label_counts)
            
            if nb_samples >= 50 and nb_classes > 3 and min_samples_per_class >= 3:
                # Cross-validation normale pour datasets suffisamment grands
                self.cv_scores = cross_val_score(self.model, self.vectorizer.transform(processed_texts), labels, cv=StratifiedKFold(n_splits=5, shuffle=True, random_state=42))
            elif nb_samples >= 20 and min_samples_per_class >= 2:
                # Cross-validation légère pour datasets moyens
                from sklearn.model_selection import KFold
                n_splits = min(3, min_samples_per_class, nb_samples // 5)
                n_splits = max(2, n_splits)  # Au minimum 2 splits
                self.cv_scores = cross_val_score(self.model, self.vectorizer.transform(processed_texts), labels, cv=KFold(n_splits=n_splits, shuffle=True, random_state=42))
            else:
                # Pour les très petits datasets, utiliser l'accuracy simple
                print("Dataset très petit, utilisation de l'accuracy simple sans cross-validation")
                self.cv_scores = np.array([self.accuracy])
        except Exception as e:
            print(f"Erreur cross-validation: {e}. Utilisation de l'accuracy simple.")
            self.cv_scores = np.array([self.accuracy])
        
        print(f"Cross-validation accuracy: {self.cv_scores.mean():.3f} (+/- {self.cv_scores.std() * 2:.3f})")

        return self.accuracy, self.report, self.cm, y_test, y_pred

    def predict(self, cv_text: str):
        if self.model is None or self.vectorizer is None:
            raise ValueError("Modèle non entraîné")
        processed_text = self.preprocess_text(cv_text)
        text_tfidf = self.vectorizer.transform([processed_text])
        prediction = self.model.predict(text_tfidf)[0]
        probabilities = self.model.predict_proba(text_tfidf)[0]
        proba_dict = {profile: float(prob) for profile, prob in zip(self.model.classes_, probabilities)}
        confidence = float(max(probabilities))
        return prediction, confidence, proba_dict

    def save_model(self, model_path: str = 'cv_classifier_model.pkl'):
        model_data = {
            'model': self.model, 
            'vectorizer': self.vectorizer,
            'accuracy': self.accuracy,
            'report': self.report,
            'cm': self.cm,
            'cv_scores': self.cv_scores,
            'nb_samples': self.nb_samples
        }
        with open(model_path, 'wb') as f:
            pickle.dump(model_data, f)

    def load_model(self, model_path: str = 'cv_classifier_model.pkl'):
        with open(model_path, 'rb') as f:
            model_data = pickle.load(f)
        self.model = model_data['model']
        self.vectorizer = model_data['vectorizer']
        self.accuracy = model_data.get('accuracy')
        self.report = model_data.get('report')
        self.cm = model_data.get('cm')
        self.cv_scores = model_data.get('cv_scores')
        self.nb_samples = model_data.get('nb_samples')